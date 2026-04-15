"""Extract all text content from the PPTX file for comparison with the markdown."""
from pptx import Presentation

prs = Presentation(r'c:\Users\u2973735\ParallelMixer\documentation\presentation\Ventilator_Control_Architecture.pptx')

for i, slide in enumerate(prs.slides):
    print(f"\n{'='*70}")
    print(f"SLIDE {i+1}: Layout={slide.slide_layout.name}")
    print(f"{'='*70}")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                text = p.text.strip()
                if text:
                    indent = "  " * p.level if p.level else ""
                    bold_marker = "[B] " if (p.runs and p.runs[0].font.bold) else ""
                    print(f"  {indent}{bold_marker}{text}")
