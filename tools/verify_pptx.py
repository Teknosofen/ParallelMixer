from pptx import Presentation
p = Presentation(r'c:\Users\u2973735\ParallelMixer\documentation\presentation\Ventilator_Control_Architecture.pptx')
print(f'Slides: {len(p.slides)}')
for i, s in enumerate(p.slides):
    title = s.shapes.title.text if s.shapes.title else '(no title)'
    print(f'  {i+1}. [{s.slide_layout.name}] {title}')
