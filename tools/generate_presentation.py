"""
Generate a Ventilator Control Architecture presentation using the Lund University template.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import os

TEMPLATE = r'c:\Users\u2973735\ParallelMixer\documentation\presentation\LundaLogger_2026-02-27 -3rd-try.pptx'
OUTPUT   = r'c:\Users\u2973735\ParallelMixer\documentation\presentation\Ventilator_Control_Architecture.pptx'

prs = Presentation(TEMPLATE)

# Delete existing slides (iterate in reverse)
xml_slides = prs.slides._sldIdLst
for sldId in list(xml_slides):
    xml_slides.remove(sldId)

# --- Layout indices (from inspection) ---
LY_TITLE       = 0   # 1.1 - Start Blue graphic  (ph0=title, ph1=subtitle, ph13=body)
LY_SECTION_B   = 3   # 2.1 - Section divider blue (ph0=title, ph1=body)
LY_CONTENT_W   = 5   # 3.1 - Heading+bullet white (ph0=title, ph14=subtitle, ph18=content)
LY_CONTENT_S   = 6   # 3.2 - Heading+bullet snow  (ph0=title, ph14=subtitle, ph18=content)
LY_TWO_COL     = 9   # 3.6 - 2 column snow        (ph0=title, ph14=subtitle, ph19=left, ph20=right)
LY_END         = 28  # 4.2 - End logo on graphic

layouts = prs.slide_layouts

# ─── helpers ─────────────────────────────────────────────────────────

def add_slide(layout_idx):
    return prs.slides.add_slide(layouts[layout_idx])

def set_ph(slide, idx, text):
    """Set text of a placeholder by its idx."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            ph.text = text
            return ph
    return None

def set_ph_rich(slide, idx, paragraphs):
    """Set placeholder content with structured bullet paragraphs.
    paragraphs: list of (text, level) or (text, level, bold) tuples.
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            tf.clear()
            for i, item in enumerate(paragraphs):
                if len(item) == 3:
                    text, level, bold = item
                else:
                    text, level = item
                    bold = False
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = text
                p.level = level
                if bold and p.runs:
                    p.runs[0].font.bold = True
            return ph
    return None

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1 – Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_TITLE)
set_ph(s, 0, "Three-Layer Ventilator\nControl Architecture")
set_ph(s, 1, "Medical Ventilator System Design")
set_ph(s, 13, "Version 1.0 – 2026-02-01")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2 – Overview
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_W)
set_ph(s, 0, "Architecture Overview")
set_ph(s, 14, "Indirect control through setpoint manipulation")
set_ph_rich(s, 18, [
    ("Three-layer control architecture", 0, True),
    ("Layer 1 – Low-Level PID Controllers (Actuator Layer)", 1),
    ("Precise, linearized control of individual actuators", 2),
    ("Layer 2 – Low-Level Controller (LLC) – Breathing Phase Control", 1),
    ("Manages insufflation / exsufflation cycle", 2),
    ("Layer 3 – High-Level Controller (HLC) – Ventilation Mode Control", 1),
    ("8-state cycle implementing complete ventilation modes", 2),
    ("", 0),
    ("Key principle", 0, True),
    ("Higher layers modify setpoints → lower layers react automatically", 1),
    ("No explicit commands between layers", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3 – Section: Layer 1
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_SECTION_B)
set_ph(s, 0, "Layer 1")
set_ph(s, 1, "Low-Level PID Controllers\n(Actuator Layer)")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4 – Layer 1: Controllers
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_S)
set_ph(s, 0, "Layer 1 – PID Controllers")
set_ph(s, 14, "Four independent PID controllers with sensor scaling")
set_ph_rich(s, 18, [
    ("O₂ Flow Controller", 0, True),
    ("Controls oxygen flow rate with sensor linearization", 1),
    ("Receives flow setpoint from LLC", 1),
    ("Air Flow Controller", 0, True),
    ("Controls air flow rate with sensor linearization", 1),
    ("Receives flow setpoint from LLC", 1),
    ("Expiratory Valve Controller", 0, True),
    ("Controls expiratory pressure / valve position", 1),
    ("Receives pressure setpoint from LLC", 1),
    ("Blower Flow Controller", 0, True),
    ("Controls blower flow rate with sensor linearization", 1),
    ("Receives flow setpoint from LLC", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5 – Layer 1: Characteristics
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_W)
set_ph(s, 0, "Layer 1 – Characteristics")
set_ph(s, 14, "Fast, independent control loops")
set_ph_rich(s, 18, [
    ("Each controller operates independently", 0),
    ("Handles non-linearities in sensors and actuators", 0),
    ("Provides fast, stable control loops (100–1000 Hz)", 0),
    ("No knowledge of breathing phases or ventilation modes", 0),
    ("", 0),
    ("Responsibilities", 0, True),
    ("Sensor scaling and compensation", 1),
    ("Actuator linearization", 1),
    ("Setpoint tracking with PID (or more advanced) control", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6 – Section: Layer 2
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_SECTION_B)
set_ph(s, 0, "Layer 2")
set_ph(s, 1, "Low-Level Controller (LLC)\nBreathing Phase Control")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7 – LLC: Two States (two-column)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_TWO_COL)
set_ph(s, 0, "LLC – Breathing Phase States")
set_ph(s, 14, "Two alternating states controlled by setpoints and triggers")
set_ph_rich(s, 19, [
    ("Insufflation (Inspiration)", 0, True),
    ("Delivers pressure/flow to patient", 1),
    ("Controls gas via O₂ and Air flow controllers", 1),
    ("Maintains target inspiratory pressure", 1),
    ("Pressure limiting: reduces flow when pressure approaches limit", 1),
])
set_ph_rich(s, 20, [
    ("Exsufflation (Expiration)", 0, True),
    ("Releases pressure/flow from patient", 1),
    ("Controls expiratory valve for pressure release", 1),
    ("Maintains PEEP (baseline pressure)", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8 – LLC: Parameters
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_W)
set_ph(s, 0, "LLC – Parameters (Set by HLC)")
set_ph(s, 14, "All parameters written by HLC – no explicit commands")
set_ph_rich(s, 18, [
    ("FiO₂ – Fraction of inspired oxygen (0.21 – 1.0)", 0),
    ("Insufflation Pressure – Target during insufflation", 0),
    ("Pressure Limit – Maximum allowed airway pressure", 0),
    ("Exsufflation Pressure (PEEP) – Target during exsufflation", 0),
    ("Insufflation Flow – Target flow rate during insufflation", 0),
    ("Bias Flow – Continuous flow during exsufflation", 0),
    ("Max Inspiratory Flow – Maximum allowed inspiratory flow", 0),
    ("Flow Trigger Threshold – Detects patient effort (flow)", 0),
    ("Pressure Trigger Threshold – Detects patient effort (pressure)", 0),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9 – LLC: Transitions & Communication
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_TWO_COL)
set_ph(s, 0, "LLC – Transitions & Communication")
set_ph(s, 14, "Purely reactive controller – no explicit commands from HLC")
set_ph_rich(s, 19, [
    ("State Transitions", 0, True),
    ("Patient Effort Detection", 1, True),
    ("Monitors flow & pressure vs trigger thresholds", 2),
    ("Inspiratory effort → triggers insufflation", 2),
    ("Expiratory effort → triggers exsufflation", 2),
    ("Setpoint-Induced Transitions", 1, True),
    ("HLC modifies setpoints → trigger criteria fulfilled", 2),
    ("Example: HLC sets PEEP=100 cmH₂O → LLC shifts to exsufflation", 2),
])
set_ph_rich(s, 20, [
    ("Signals to HLC", 0, True),
    ("State change events (insuff ↔ exsuff)", 1),
    ("Measured tidal volume (integrated flow)", 1),
    ("Pressure limit events", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10 – Section: Layer 3
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_SECTION_B)
set_ph(s, 0, "Layer 3")
set_ph(s, 1, "High-Level Controller (HLC)\nVentilation Mode Control")

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11 – HLC: 8-State Machine Overview
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_S)
set_ph(s, 0, "HLC – 8-State Sequential Machine")
set_ph(s, 14, "Two phases × four states per phase")
set_ph_rich(s, 18, [
    ("Inspiration Sequence", 0, True),
    ("1. Non-triggerable – Fixed inspiration, triggers ignored", 1),
    ("2. Support – Patient can trigger pressure support", 1),
    ("3. Synch – Patient-synchronized transition (default: 0)", 1),
    ("4. Pause – Inspiratory hold for plateau measurement", 1),
    ("", 0),
    ("Expiration Sequence", 0, True),
    ("5. Non-triggerable – Fixed expiration, triggers ignored", 1),
    ("6. Support – Patient can trigger expiratory support (rare)", 1),
    ("7. Synch – Patient-synchronized transition (default: 0)", 1),
    ("8. Pause – Expiratory hold for measurements", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12 – HLC: Inspiration States Detail
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_W)
set_ph(s, 0, "HLC – Inspiration States")
set_ph(s, 14, "States 1–4 of the ventilation cycle")
set_ph_rich(s, 18, [
    ("1. Non-triggerable", 0, True),
    ("Delivers mandatory breath", 1),
    ("Patient triggers are ignored", 1),
    ("Duration set by timing parameters", 1),
    ("2. Support", 0, True),
    ("Patient can trigger additional inspiratory support pressure", 1),
    ("Allows patient to augment mechanical breath", 1),
    ("3. Synch (Synchronization)", 0, True),
    ("LLC exsufflation signal → HLC advances to Pause", 1),
    ("Allows patient to terminate inspiration", 1),
    ("4. Pause (Inspiratory Plateau)", 0, True),
    ("Holds delivered volume – no gas added or released", 1),
    ("Used for plateau pressure measurement", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 13 – HLC: Expiration States Detail
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_W)
set_ph(s, 0, "HLC – Expiration States")
set_ph(s, 14, "States 5–8 of the ventilation cycle")
set_ph_rich(s, 18, [
    ("5. Non-triggerable", 0, True),
    ("Ensures minimum expiration time", 1),
    ("Patient triggers are ignored", 1),
    ("6. Support", 0, True),
    ("Patient can trigger expiratory support pressure", 1),
    ("Rare in typical ventilation modes", 1),
    ("7. Synch (Synchronization)", 0, True),
    ("LLC insufflation signal → HLC advances to Pause", 1),
    ("Allows patient to initiate next breath", 1),
    ("8. Pause (Expiratory)", 0, True),
    ("Maintains expiratory state", 1),
    ("Used for breath holds or measurements", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 14 – State Transition Logic
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_W)
set_ph(s, 0, "HLC – State Transition Logic")
set_ph(s, 14, "Simple, deterministic transitions – no parameter interpretation")
set_ph_rich(s, 18, [
    ("States advance based on:", 0, True),
    ("Time elapsed – each state has a pre-calculated absolute duration", 1),
    ("Zero-duration skipping – states with time = 0 are immediately skipped", 1),
    ("Patient synchronization – during Synch states, LLC signals trigger advance", 1),
    ("Insp Synch: LLC exsufflation signal → advance to Insp Pause", 2),
    ("Exp Synch: LLC insufflation signal → advance to Exp Pause", 2),
    ("", 0),
    ("Key design principle", 0, True),
    ("The HLC operates exclusively on absolute times (ms/s)", 1),
    ("It does NOT interpret RR, I:E ratio, pause %, or any user parameters", 1),
    ("All parameter-to-time conversion is done by the UI Layer above", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 15 – HLC Absolute Timing Parameters
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_S)
set_ph(s, 0, "HLC – Timing Interface")
set_ph(s, 14, "Eight absolute durations – the only timing inputs to the HLC")
set_ph_rich(s, 18, [
    ("Inspiration phase", 0, True),
    ("Time_Insp_NonTrig – Mandatory inspiration duration", 1),
    ("Time_Insp_Support – Inspiration support window", 1),
    ("Time_Insp_Synch – Max wait for patient-synchronized termination", 1),
    ("Time_Insp_Pause – Inspiratory pause / plateau hold", 1),
    ("", 0),
    ("Expiration phase", 0, True),
    ("Time_Exp_NonTrig – Mandatory expiration duration", 1),
    ("Time_Exp_Support – Expiration support window", 1),
    ("Time_Exp_Synch – Max wait for patient-synchronized termination", 1),
    ("Time_Exp_Pause – Expiratory pause / hold", 1),
    ("", 0),
    ("Sum of all 8 durations = one complete breathing cycle", 0, True),
    ("Active Insp. Time = NonTrig + Support + Synch (pause excluded for Vt calc)", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 16 – User Interface Layer
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_TWO_COL)
set_ph(s, 0, "User Interface Layer")
set_ph(s, 14, "Translates clinician parameters into HLC absolute times")
set_ph_rich(s, 19, [
    ("UI Layer accepts user parameters", 0, True),
    ("Respiratory Rate (RR)", 1),
    ("I:E Ratio", 1),
    ("Inspiratory Pause %", 1),
    ("Expiratory Pause %", 1),
    ("Non-triggerable / Support / Synch fractions", 1),
    ("Tidal Volume, FiO₂, pressures, etc.", 1),
])
set_ph_rich(s, 20, [
    ("UI Layer solves the equations", 0, True),
    ("Cycle Time = 60 / RR", 1),
    ("Insp Time = Cycle × I/(I+E)", 1),
    ("Exp Time = Cycle × E/(I+E)", 1),
    ("Distributes durations across 8 states", 1),
    ("Writes absolute times + setpoints to HLC", 1),
    ("", 0),
    ("Benefit: different UIs (touchscreen, remote,", 0, True),
    ("automated protocols) share the same", 0),
    ("deterministic HLC interface", 0),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 17 – Tidal Volume Control
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_S)
set_ph(s, 0, "Tidal Volume Control")
set_ph(s, 14, "Breath-to-breath adaptive control with pressure safety")
set_ph_rich(s, 18, [
    ("Initial Flow Calculation", 0, True),
    ("Insufflation Flow = Vt / Active Inspiration Time", 1),
    ("Active Insp. Time = Time(Non-trig) + Time(Support) + Time(Synch)", 1),
    ("", 0),
    ("Breath-to-Breath Adaptation", 0, True),
    ("At each inspi→expi transition:", 1),
    ("Measure actual Vt_measured", 2),
    ("Calculate ΔVt = Vt_target − Vt_measured", 2),
    ("Gradually adjust insufflation flow setpoint", 2),
    ("Prevents oscillations, compensates for compliance changes", 2),
    ("", 0),
    ("Pressure-Limited Volume Control", 0, True),
    ("LLC pressure controllers auto-reduce flow near Pressure Limit", 1),
    ("Vt_measured < Vt_target → volume sacrificed for safety", 1),
    ("Two modes: Volume-controlled  /  Pressure-limited volume-controlled", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 16 – Pressure Support Mechanism
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_W)
set_ph(s, 0, "Pressure Support Mechanism")
set_ph(s, 14, "Rewarding patient effort through pressure boosts")
set_ph_rich(s, 18, [
    ("Inspiration Support State", 0, True),
    ("When patient trigger detected during Insp Support state:", 1),
    ("LLC Insuff Pressure = Inspiratory Pressure + Insp Support Pressure", 1),
    ("Provides additional pressure boost rewarding inspiratory effort", 1),
    ("Common in PSV (Pressure Support Ventilation) modes", 1),
    ("", 0),
    ("Expiration Support State", 0, True),
    ("When patient trigger detected during Exp Support state:", 1),
    ("LLC Insuff Pressure = PEEP + Exp Support Pressure", 1),
    ("Less common – available for specific clinical scenarios", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 17 – Inspiratory Pause Implementation
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_S)
set_ph(s, 0, "Inspiratory Pause Implementation")
set_ph(s, 14, "Elegant setpoint manipulation – no special pause logic in LLC")
set_ph_rich(s, 18, [
    ("Purpose", 0, True),
    ("Hold delivered volume for plateau pressure measurement", 1),
    ("Assess static lung compliance", 1),
    ("Equilibrate airway and alveolar pressures", 1),
    ("", 0),
    ("Mechanism (setpoint manipulation)", 0, True),
    ("Bias Flow → 0  (no gas added by insufflation controllers)", 1),
    ("PEEP setpoint → 100 cmH₂O  (forces LLC into exsufflation mode)", 1),
    ("Actual pressure ≪ 100 → expiratory valve stays closed", 1),
    ("Result: system is mechanically \"locked\"", 1),
    ("No gas flows in, no gas escapes", 2),
    ("Volume held constant", 2),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 18 – HLC Control: State-Dependent Setpoints
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_W)
set_ph(s, 0, "HLC – Indirect Control Philosophy")
set_ph(s, 14, "No explicit commands – only setpoint manipulation")
set_ph_rich(s, 18, [
    ("HLC controls LLC by:", 0, True),
    ("Setting insufflation pressure based on state & triggers", 1),
    ("Setting pressure limit for safety", 1),
    ("Setting insufflation flow (Vt, timing, adaptation)", 1),
    ("Adjusting bias flow & PEEP for pause phases", 1),
    ("Adjusting trigger thresholds per HLC state", 1),
    ("", 0),
    ("State-Dependent Setpoints (key states)", 0, True),
    ("Insp Non-trig: Insp Pressure | Normal PEEP | Normal Bias", 1),
    ("Insp Support (triggered): Insp Press + Support | Normal PEEP", 1),
    ("Insp Pause: N/A | PEEP=100 cmH₂O | Bias=0", 1),
    ("Exp states: N/A or PEEP+Support | Normal PEEP | Normal Bias", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 19 – Communication Between Layers
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_S)
set_ph(s, 0, "Communication Between Layers")
set_ph(s, 14, "Clean interfaces with minimal coupling")
set_ph_rich(s, 18, [
    ("HLC → LLC  (Setpoints)", 0, True),
    ("Pressure setpoints (insp, PEEP, limit)", 1),
    ("Flow setpoints (insufflation, bias, max insp)", 1),
    ("Trigger thresholds (flow & pressure)", 1),
    ("", 0),
    ("LLC → HLC  (Signals)", 0, True),
    ("State change events (insufflation ↔ exsufflation)", 1),
    ("Measured tidal volume (integrated flow)", 1),
    ("Pressure limit events", 1),
    ("", 0),
    ("LLC → Layer 1  (Setpoints)", 0, True),
    ("Individual flow controller setpoints", 1),
    ("Expiratory valve pressure setpoint", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 20 – Example Ventilation Modes (two-column)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_TWO_COL)
set_ph(s, 0, "Example Ventilation Modes")
set_ph(s, 14, "Various clinical modes via timing & parameter configuration")
set_ph_rich(s, 19, [
    ("Volume Control (VC)", 0, True),
    ("Insp Non-trig delivers full breath", 1),
    ("Support & Synch skipped (time=0)", 1),
    ("Time-triggered, volume-cycled", 1),
    ("", 0),
    ("Assist Control (AC)", 0, True),
    ("Backup mandatory breath timing", 1),
    ("Exp Synch active → patient triggering", 1),
    ("Full support on every breath", 1),
])
set_ph_rich(s, 20, [
    ("Pressure Support (PSV)", 0, True),
    ("Insp Support active with pressure", 1),
    ("Insp & Exp Synch active", 1),
    ("Fully patient-triggered & cycled", 1),
    ("", 0),
    ("SIMV", 0, True),
    ("Mandatory (Non-trig) + Supported breaths", 1),
    ("Synch states coordinate with patient", 1),
    ("Combines mandatory & spontaneous", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 21 – Safety & Design Advantages (two-column)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_TWO_COL)
set_ph(s, 0, "Safety & Design Advantages")
set_ph(s, 14, "")
set_ph_rich(s, 19, [
    ("Safety Considerations", 0, True),
    ("Pressure limiting always active (barotrauma prevention)", 1),
    ("Non-trig states ensure minimum ventilation", 1),
    ("Gradual adaptation prevents sudden volume changes", 1),
    ("Setpoint sanity checks in HLC", 1),
    ("Watchdog timers per state", 1),
])
set_ph_rich(s, 20, [
    ("Design Advantages", 0, True),
    ("Modularity – clear layer responsibilities", 1),
    ("No direct commands – reduced coupling", 1),
    ("Extensibility – new modes via config only", 1),
    ("Patient synchronization built-in", 1),
    ("Breath-to-breath adaptability", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 22 – Implementation Notes
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_CONTENT_W)
set_ph(s, 0, "Implementation Notes")
set_ph(s, 14, "Practical considerations for system design")
set_ph_rich(s, 18, [
    ("Execution Rates", 0, True),
    ("Layer 1 PID loops: 100 – 1000 Hz", 1),
    ("LLC: event-driven architecture", 1),
    ("HLC: fixed time base (10 – 50 Hz)", 1),
    ("", 0),
    ("Parameter Updates", 0, True),
    ("Validate all parameter changes before applying", 1),
    ("Rate-limit parameter changes to prevent instability", 1),
    ("Log all setpoint changes for debugging & clinical review", 1),
    ("", 0),
    ("Testing Strategy", 0, True),
    ("Test each layer independently first", 1),
    ("Verify LLC response to setpoint changes", 1),
    ("Verify HLC state sequences for all modes", 1),
    ("Patient simulator testing with varying compliance/resistance", 1),
])

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 23 – End
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
s = add_slide(LY_END)

# ─── Save ────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Presentation saved to: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
