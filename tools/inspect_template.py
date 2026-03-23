from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r'c:\Users\u2973735\ParallelMixer\documentation\presentation\LundaLogger_2026-02-27 -3rd-try.pptx')

print('=== Slide dimensions ===')
print(f'Width: {prs.slide_width}, Height: {prs.slide_height}')
print(f'Width (inches): {prs.slide_width / 914400}, Height (inches): {prs.slide_height / 914400}')

print('\n=== Slide Layouts ===')
for i, layout in enumerate(prs.slide_layouts):
    print(f'  Layout {i}: {layout.name}')
    for ph in layout.placeholders:
        print(f'    Placeholder {ph.placeholder_format.idx}: {ph.name} (type={ph.placeholder_format.type})')

print('\n=== Existing Slides ===')
for i, slide in enumerate(prs.slides):
    print(f'\nSlide {i+1}: Layout={slide.slide_layout.name}')
    for shape in slide.shapes:
        print(f'  Shape: {shape.shape_type}, name={shape.name}, pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})')
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                text = p.text[:100] if p.text else ''
                if p.runs:
                    r = p.runs[0]
                    try:
                        clr = str(r.font.color.rgb) if r.font.color and r.font.color.type else None
                    except:
                        clr = 'inherited'
                    font_info = f'font={r.font.name}, size={r.font.size}, bold={r.font.bold}, color={clr}'
                else:
                    font_info = 'no runs'
                print(f'    Para: "{text}" [{font_info}]')
